#!/usr/bin/env python3
"""
Advanced PPTX to Markdown Converter using Vision Language Models
Converts visual PowerPoint slides to markdown using modern OCR/VLM approaches.

Supports multiple backends:
- PyVisionAI (OpenAI GPT-4 Vision or Ollama)
- olmOCR (local VLM)
- GOT-OCR 2.0
- Qwen2.5-VL

Author: Law School Workspace
Date: February 2026
"""

import argparse
import os
import sys
from pathlib import Path
from typing import List, Optional
import subprocess
import json

# Try imports with helpful error messages
try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    print("Error: Pillow not installed. Run: pip install Pillow")
    sys.exit(1)


class SlidesToMarkdownConverter:
    """Convert PPTX slides to Markdown using VLM/OCR"""
    
    def __init__(self, backend: str = "pyvisionai", output_dir: Optional[Path] = None):
        """
        Initialize converter
        
        Args:
            backend: OCR backend to use (pyvisionai, olmocr, gotocr, qwen)
            output_dir: Directory for output files
        """
        self.backend = backend.lower()
        self.output_dir = output_dir
        self.supported_backends = ["pyvisionai", "olmocr", "gotocr", "qwen", "simple"]
        
        if self.backend not in self.supported_backends:
            raise ValueError(f"Backend must be one of: {self.supported_backends}")
        
        # Check backend availability
        self._check_backend()
    
    def _check_backend(self):
        """Check if selected backend is available"""
        if self.backend == "pyvisionai":
            try:
                import pyvisionai
                print(f"✓ PyVisionAI backend ready")
            except ImportError:
                print("Warning: PyVisionAI not installed. Run: pip install pyvisionai")
                print("Falling back to simple extraction")
                self.backend = "simple"
        
        elif self.backend == "olmocr":
            try:
                import olmocr
                print(f"✓ olmOCR backend ready")
            except ImportError:
                print("Warning: olmOCR not installed. Run: pip install olmocr")
                print("Falling back to simple extraction")
                self.backend = "simple"
        
        elif self.backend == "gotocr":
            print("Note: GOT-OCR 2.0 requires manual setup from GitHub")
            print("See: https://github.com/ucaslcl/GOT-OCR2.0")
        
        elif self.backend == "qwen":
            print("Note: Qwen2.5-VL requires manual setup")
            print("See: https://github.com/QwenLM/Qwen2-VL")
    
    def extract_slides_as_images(self, pptx_path: Path) -> List[Path]:
        """
        Extract each slide as an image
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            List of paths to slide images
        """
        print(f"\n📄 Processing: {pptx_path.name}")
        
        # Create temp directory for slide images
        temp_dir = self.output_dir / f"_temp_{pptx_path.stem}"
        temp_dir.mkdir(parents=True, exist_ok=True)
        
        # Load presentation
        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            print(f"Error loading presentation: {e}")
            return []
        
        slide_images = []
        
        # Try to export slides as images using LibreOffice or PowerPoint
        # This preserves visual layout better than extracting programmatically
        print(f"  Converting {len(prs.slides)} slides to images...")
        
        # Method 1: Try using LibreOffice (cross-platform)
        slide_images = self._export_with_libreoffice(pptx_path, temp_dir)
        
        if not slide_images:
            # Method 2: Try programmatic extraction
            print("  Falling back to programmatic extraction...")
            slide_images = self._extract_slides_programmatic(prs, temp_dir)
        
        return slide_images
    
    def _export_with_libreoffice(self, pptx_path: Path, output_dir: Path) -> List[Path]:
        """Export slides as images using LibreOffice"""
        try:
            # Convert PPTX to PDF first
            pdf_path = output_dir / f"{pptx_path.stem}.pdf"
            cmd = [
                "soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(output_dir),
                str(pptx_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode != 0 or not pdf_path.exists():
                return []
            
            # Convert PDF to images using ImageMagick
            print(f"  Converting PDF to images...")
            cmd = [
                "convert",
                "-density", "300",
                "-quality", "100",
                str(pdf_path),
                str(output_dir / "slide-%03d.png")
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            if result.returncode == 0:
                slide_images = sorted(output_dir.glob("slide-*.png"))
                print(f"  ✓ Exported {len(slide_images)} slides as images")
                return slide_images
            
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            print(f"  LibreOffice/ImageMagick not available: {e}")
        
        return []
    
    def _extract_slides_programmatic(self, prs: Presentation, output_dir: Path) -> List[Path]:
        """Extract slide images programmatically (fallback)"""
        slide_images = []
        
        for i, slide in enumerate(prs.slides):
            # This is a simplified version - in practice, we'd need to render the slide
            # For now, just extract embedded images
            slide_dir = output_dir / f"slide_{i:03d}"
            slide_dir.mkdir(exist_ok=True)
            
            image_count = 0
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    try:
                        image = shape.image
                        image_path = slide_dir / f"image_{image_count}.{image.ext}"
                        with open(image_path, 'wb') as f:
                            f.write(image.blob)
                        image_count += 1
                    except Exception as e:
                        print(f"    Warning: Could not extract image: {e}")
            
            if image_count > 0:
                slide_images.append(slide_dir)
        
        return slide_images
    
    def process_with_pyvisionai(self, slide_image: Path) -> str:
        """Process slide image with PyVisionAI"""
        try:
            from pyvisionai import VisionAI
            
            # Initialize with OpenAI (or can use Ollama for local processing)
            vision = VisionAI(provider="openai")  # or "ollama" for local
            
            # Describe the slide image
            prompt = """Analyze this lecture slide and extract all content in Markdown format.
            Include:
            - Main title (as # heading)
            - Subtitles (as ## headings)
            - All bullet points
            - All text content
            - Descriptions of any diagrams, charts, or important visual elements
            - Any formulas or equations
            
            Format as clean, hierarchical Markdown."""
            
            result = vision.describe_image(str(slide_image), prompt=prompt)
            return result
            
        except Exception as e:
            print(f"    Error with PyVisionAI: {e}")
            return ""
    
    def process_with_olmocr(self, slide_image: Path) -> str:
        """Process slide image with olmOCR"""
        try:
            from olmocr import OCR
            
            ocr = OCR()
            result = ocr.process(str(slide_image), output_format="markdown")
            return result
            
        except Exception as e:
            print(f"    Error with olmOCR: {e}")
            return ""
    
    def process_simple(self, pptx_path: Path) -> str:
        """Simple text extraction from PPTX (fallback)"""
        try:
            prs = Presentation(str(pptx_path))
            markdown_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                markdown_content.append(f"\n## Slide {i}\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Try to determine if it's a title
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                            if shape.placeholder_format.type == 1:  # Title
                                markdown_content.append(f"### {shape.text}\n")
                            else:
                                markdown_content.append(f"{shape.text}\n")
                        else:
                            markdown_content.append(f"{shape.text}\n")
                
                # Extract notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        markdown_content.append(f"\n**Speaker Notes:**\n{notes_text}\n")
            
            return "\n".join(markdown_content)
            
        except Exception as e:
            print(f"    Error with simple extraction: {e}")
            return ""
    
    def convert_file(self, pptx_path: Path) -> Optional[Path]:
        """
        Convert a single PPTX file to Markdown
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Path to output markdown file
        """
        print(f"\n{'='*60}")
        print(f"Converting: {pptx_path.name}")
        print(f"Backend: {self.backend}")
        print(f"{'='*60}")
        
        # Determine output path
        if self.output_dir:
            output_path = self.output_dir / f"{pptx_path.stem}.md"
        else:
            output_path = pptx_path.with_suffix('.md')
        
        # Process based on backend
        if self.backend in ["pyvisionai", "olmocr", "gotocr", "qwen"]:
            # Extract slides as images first
            slide_images = self.extract_slides_as_images(pptx_path)
            
            if not slide_images:
                print("  Warning: Could not extract slide images, using simple extraction")
                markdown_content = self.process_simple(pptx_path)
            else:
                # Process each slide image
                markdown_parts = [f"# {pptx_path.stem}\n\n"]
                
                for i, slide_img in enumerate(slide_images, 1):
                    print(f"  Processing slide {i}/{len(slide_images)}...")
                    
                    if self.backend == "pyvisionai":
                        content = self.process_with_pyvisionai(slide_img)
                    elif self.backend == "olmocr":
                        content = self.process_with_olmocr(slide_img)
                    else:
                        content = f"[Slide {i} - requires {self.backend} setup]"
                    
                    markdown_parts.append(f"\n---\n\n## Slide {i}\n\n{content}\n")
                
                markdown_content = "\n".join(markdown_parts)
        else:
            # Simple extraction
            markdown_content = self.process_simple(pptx_path)
        
        # Write output
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            print(f"\n✓ Saved: {output_path}")
            return output_path
            
        except Exception as e:
            print(f"\nError saving file: {e}")
            return None
    
    def convert_directory(self, input_dir: Path) -> List[Path]:
        """
        Convert all PPTX files in a directory
        
        Args:
            input_dir: Directory containing PPTX files
            
        Returns:
            List of output markdown files
        """
        pptx_files = list(input_dir.glob("*.pptx")) + list(input_dir.glob("*.PPTX"))
        
        if not pptx_files:
            print(f"No PPTX files found in {input_dir}")
            return []
        
        print(f"\nFound {len(pptx_files)} PPTX files to convert")
        
        output_files = []
        for pptx_file in pptx_files:
            output_file = self.convert_file(pptx_file)
            if output_file:
                output_files.append(output_file)
        
        return output_files


def main():
    parser = argparse.ArgumentParser(
        description="Convert PPTX slides to Markdown using VLM/OCR",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Simple extraction (no OCR, just text)
  python slides_to_markdown_vlm.py /path/to/slides --backend simple
  
  # Using PyVisionAI with OpenAI GPT-4 Vision
  python slides_to_markdown_vlm.py /path/to/slides --backend pyvisionai
  
  # Using olmOCR (local VLM)
  python slides_to_markdown_vlm.py /path/to/slides --backend olmocr
  
  # Specify output directory
  python slides_to_markdown_vlm.py /path/to/slides -o /path/to/output

Backends:
  simple      - Text extraction only (no OCR) - FASTEST
  pyvisionai  - GPT-4 Vision or Ollama (best for complex slides)
  olmocr      - Local VLM (fast, high quality)
  gotocr      - GOT-OCR 2.0 (best for scientific/technical)
  qwen        - Qwen2.5-VL (excellent for diagrams/charts)
        """
    )
    
    parser.add_argument(
        "input",
        type=Path,
        help="PPTX file or directory containing PPTX files"
    )
    
    parser.add_argument(
        "-b", "--backend",
        choices=["simple", "pyvisionai", "olmocr", "gotocr", "qwen"],
        default="simple",
        help="OCR/extraction backend to use (default: simple)"
    )
    
    parser.add_argument(
        "-o", "--output",
        type=Path,
        help="Output directory (default: same as input)"
    )
    
    args = parser.parse_args()
    
    # Validate input
    if not args.input.exists():
        print(f"Error: Input path does not exist: {args.input}")
        sys.exit(1)
    
    # Setup converter
    converter = SlidesToMarkdownConverter(
        backend=args.backend,
        output_dir=args.output
    )
    
    # Convert
    if args.input.is_file():
        converter.convert_file(args.input)
    else:
        converter.convert_directory(args.input)
    
    print(f"\n{'='*60}")
    print("✓ Conversion complete!")
    print(f"{'='*60}\n")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""
PPTX to text converter.
Usage: python3 pptx_to_text.py <pptx_file> [output_file]
"""

import sys
from pptx import Presentation

def pptx_to_text(pptx_path, output_path=None):
    """Extract text from PPTX file."""
    
    if output_path is None:
        output_path = pptx_path.rsplit('.', 1)[0] + '_text.txt'
    
    prs = Presentation(pptx_path)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        for slide_number, slide in enumerate(prs.slides, 1):
            f.write(f"\n{'='*80}\n")
            f.write(f"SLIDE {slide_number}\n")
            f.write(f"{'='*80}\n\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f.write(shape.text + "\n")
    
    print(f"Text extracted to: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 pptx_to_text.py <pptx_file> [output_file]")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    pptx_to_text(pptx_path, output_path)